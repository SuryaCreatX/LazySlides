import collections 
import collections.abc
import glob
from pptx import Presentation
from pptx.util import Inches,Pt

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

class text2s:
    def __init__(self,make):
        self.make = make
        for path in glob.glob("textfile/*.txt"):
            f = open(path, "r")
            a=f.read()
            b=[]
            i=0

            while(i<len(a)):
                slice1=a[i:i+100]+"\n"+a[i+100:i+200]+"\n"+a[i+200:i+300]+"\n"+a[i+300:i+400]+"\n"+a[i+400:i+500]+"\n"+a[i+500:i+600]+"\n"+a[i+600:i+700]+"\n"+a[i+700:i+800]+"\n"+a[i+800:i+900]+"\n"+a[i+900:i+1000]+"\n"+a[i+1000:i+1100]+"\n"+a[i+1100:i+1200]+"\n"+a[i+1200:i+1300]+"\n"
                b.append(slice1)
                i+=1300 

            for o in b:
                
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                top = Inches(1)
                left1 = Inches(12)
                height1 = width1 = Inches(3.5)

                txBox = slide.shapes.add_textbox(left=Inches(1), top=top, width=left1, height=height1)
                tf = txBox.text_frame
                tf.text = o
                prs.save('output/text.pptx')