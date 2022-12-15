import cv2
import pytesseract 
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

pytesseract.pytesseract.tesseract_cmd = '"C:\Program Files\Tesseract-OCR/tesseract.exe"'

class pptx1:
    def __init__(self,make):
        self.make = make

    def convert_screenshots_to_pptx(self,img2,path):
        mytext=""
        i=0
        k=[]
        texts =  pytesseract.image_to_data(img2)  
        img_path: str = path
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        top = Inches(1)
        left1 = Inches(12)
        height1 = width1 = Inches(3.5)
        for b,o in enumerate(texts.splitlines()):
            g=o.split()
            k.append(g)
        g.append("")
        g[6]=0
        k.append(g)
        m=10000
        lst1=[1]
        lst2=[1]
        for j in k:
            if len(j)==12 and i>=1:
                x,y,w,h = int(j[6]),int(j[7]),int(j[8]),int(j[9])
                if x<m and i>=2:
                    txBox = slide.shapes.add_textbox(left=7500*min(lst1), top=6000*(sum(lst2) / len(lst2)), width=w, height=h)
                    tf = txBox.text_frame
                    tf.text = mytext
                    #time.sleep(1)
                    mytext=""
                    lst1.clear()
                    lst2.clear()
                mytext+=j[11]+" "
                m=x
                lst1.append(x)
                lst2.append(y)
            i+=1
        pic = slide.shapes.add_picture(img_path, left = left1, top =top, height=height1,width=width1)
        prs.save('output/'+self.make+'.pptx')